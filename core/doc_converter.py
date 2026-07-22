"""文档格式转换"""
import os
import re
import io


class DocumentConverter:
    def __init__(self):
        self._cancel = False

    def cancel(self):
        self._cancel = True

    def convert(self, input_path, output_path, progress_callback=None):
        self._cancel = False
        ext = os.path.splitext(input_path)[1].lower()
        out_ext = os.path.splitext(output_path)[1].lower()

        if progress_callback:
            progress_callback(10, "正在转换...")

        try:
            handler = self._get_handler(ext, out_ext)
            if handler is None:
                if progress_callback:
                    progress_callback(-1, f"不支持 {ext} → {out_ext} 转换")
                return False

            result = handler(input_path, output_path, progress_callback)

            if result and os.path.exists(output_path):
                if progress_callback:
                    progress_callback(100, "转换完成")
                return True
            else:
                if progress_callback:
                    progress_callback(-1, "转换失败：输出文件未生成")
                return False

        except Exception as e:
            if progress_callback:
                progress_callback(-1, f"错误: {str(e)[:200]}")
            return False

    def _get_handler(self, in_ext, out_ext):
        key = (in_ext, out_ext)
        handlers = {
            # PDF → 其他
            (".pdf", ".docx"): self._pdf_to_docx,
            (".pdf", ".doc"): self._pdf_to_docx,
            (".pdf", ".txt"): self._pdf_to_txt,
            (".pdf", ".jpg"): self._pdf_to_image,
            (".pdf", ".jpeg"): self._pdf_to_image,
            (".pdf", ".png"): self._pdf_to_image,
            (".pdf", ".html"): self._pdf_to_html,
            (".pdf", ".pptx"): self._pdf_to_pptx,
            (".pdf", ".xlsx"): self._pdf_to_xlsx,
            # Word → 其他
            (".docx", ".pdf"): self._docx_to_pdf,
            (".docx", ".txt"): self._docx_to_txt,
            (".docx", ".html"): self._docx_to_html,
            (".docx", ".jpg"): self._docx_to_image,
            (".docx", ".png"): self._docx_to_image,
            (".docx", ".pptx"): self._docx_to_pptx,
            (".docx", ".md"): self._docx_to_md,
            (".doc", ".pdf"): self._docx_to_pdf,
            (".doc", ".txt"): self._docx_to_txt,
            (".doc", ".docx"): self._doc_copy,
            (".docx", ".doc"): self._doc_copy,
            (".doc", ".md"): self._docx_to_md,
            (".doc", ".html"): self._docx_to_html,
            # WPS → 其他
            (".wps", ".docx"): self._doc_copy,
            (".wps", ".pdf"): self._docx_to_pdf,
            (".wps", ".txt"): self._docx_to_txt,
            (".docx", ".wps"): self._doc_copy,
            (".wps", ".html"): self._docx_to_html,
            (".wps", ".md"): self._docx_to_md,
            # Excel → 其他
            (".xlsx", ".pdf"): self._xlsx_to_pdf,
            (".xlsx", ".csv"): self._xlsx_to_csv,
            (".xlsx", ".txt"): self._xlsx_to_txt,
            (".xlsx", ".jpg"): self._xlsx_to_image,
            (".xlsx", ".png"): self._xlsx_to_image,
            (".xlsx", ".html"): self._xlsx_to_html,
            (".xlsx", ".md"): self._xlsx_to_md,
            (".xls", ".xlsx"): self._xls_to_xlsx,
            (".xls", ".pdf"): self._xls_to_pdf,
            (".xls", ".csv"): self._xls_to_csv,
            (".xls", ".txt"): self._xls_to_txt,
            (".xls", ".jpg"): self._xls_to_image,
            (".xls", ".png"): self._xls_to_image,
            (".xls", ".html"): self._xls_to_html,
            (".xls", ".md"): self._xls_to_md,
            (".csv", ".xlsx"): self._csv_to_xlsx,
            (".csv", ".pdf"): self._csv_to_pdf,
            (".csv", ".txt"): self._doc_copy,
            (".csv", ".html"): self._csv_to_html,
            (".csv", ".md"): self._csv_to_md,
            (".txt", ".xlsx"): self._txt_to_xlsx,
            (".txt", ".docx"): self._txt_to_docx,
            (".txt", ".pptx"): self._txt_to_pptx,
            (".txt", ".pdf"): self._txt_to_pdf,
            (".txt", ".html"): self._txt_to_html,
            (".txt", ".md"): self._txt_to_md,
            # PPT → 其他
            (".pptx", ".pdf"): self._pptx_to_pdf,
            (".pptx", ".txt"): self._pptx_to_txt,
            (".pptx", ".jpg"): self._pptx_to_image,
            (".pptx", ".png"): self._pptx_to_image,
            (".pptx", ".docx"): self._pptx_to_docx,
            (".pptx", ".html"): self._pptx_to_html,
            (".pptx", ".md"): self._pptx_to_md,
            (".ppt", ".pptx"): self._doc_copy,
            (".ppt", ".pdf"): self._ppt_to_pdf,
            (".ppt", ".txt"): self._ppt_to_txt,
            (".pptx", ".ppt"): self._doc_copy,
            # WPS演示
            (".dps", ".pptx"): self._doc_copy,
            (".pptx", ".dps"): self._doc_copy,
            (".dps", ".pdf"): self._pptx_to_pdf,
            (".dps", ".txt"): self._pptx_to_txt,
            # WPS表格
            (".et", ".xlsx"): self._doc_copy,
            (".xlsx", ".et"): self._doc_copy,
            (".et", ".pdf"): self._xlsx_to_pdf,
            (".et", ".csv"): self._xlsx_to_csv,
            # 图片 → 文档
            (".jpg", ".pdf"): self._image_to_pdf,
            (".jpeg", ".pdf"): self._image_to_pdf,
            (".png", ".pdf"): self._image_to_pdf,
            (".bmp", ".pdf"): self._image_to_pdf,
            (".tiff", ".pdf"): self._image_to_pdf,
            (".webp", ".pdf"): self._image_to_pdf,
            (".jpg", ".docx"): self._image_to_docx,
            (".jpeg", ".docx"): self._image_to_docx,
            (".png", ".docx"): self._image_to_docx,
            (".bmp", ".docx"): self._image_to_docx,
            # HTML → 其他
            (".html", ".pdf"): self._html_to_pdf,
            (".htm", ".pdf"): self._html_to_pdf,
            (".html", ".docx"): self._html_to_docx,
            (".htm", ".docx"): self._html_to_docx,
            (".html", ".txt"): self._html_to_txt,
            (".htm", ".txt"): self._html_to_txt,
            (".html", ".md"): self._html_to_md,
            (".htm", ".md"): self._html_to_md,
            (".html", ".xlsx"): self._html_to_xlsx,
            (".htm", ".xlsx"): self._html_to_xlsx,
            # Markdown → 其他
            (".md", ".html"): self._md_to_html,
            (".md", ".pdf"): self._md_to_pdf,
            (".md", ".docx"): self._md_to_docx,
            (".md", ".txt"): self._md_to_txt,
            # EPUB → 其他
            (".epub", ".pdf"): self._epub_to_pdf,
            (".epub", ".txt"): self._epub_to_txt,
            (".epub", ".html"): self._epub_to_html,
            (".epub", ".docx"): self._epub_to_docx,
            # RTF → 其他
            (".rtf", ".txt"): self._rtf_to_txt,
            (".rtf", ".pdf"): self._rtf_to_pdf,
            (".rtf", ".docx"): self._rtf_to_docx,
            # ODT → 其他
            (".odt", ".pdf"): self._odt_to_pdf,
            (".odt", ".docx"): self._odt_to_docx,
            (".odt", ".txt"): self._odt_to_txt,
        }
        return handlers.get(key)

    # ========== 工具方法 ==========

    def _read_text(self, path):
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()

    def _write_text(self, path, text):
        with open(path, 'w', encoding='utf-8') as f:
            f.write(text)

    def _make_table_data(self, ws):
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append([str(c) if c is not None else '' for c in row])
        if not data:
            data = [['(空表格)']]
        return data

    def _safe_html(self, text):
        return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    def _build_html_page(self, body):
        return f"<html><head><meta charset='utf-8'></head><body>{body}</body></html>"

    # ========== PDF 转换 ==========

    def _pdf_to_docx(self, inp, out, cb):
        if cb: cb(20, "解析PDF...")
        from pdf2docx import Converter
        cv = Converter(inp)
        if cb: cb(50, "生成Word...")
        cv.convert(out)
        cv.close()
        return True

    def _pdf_to_txt(self, inp, out, cb):
        import fitz
        doc = fitz.open(inp)
        text = []
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text.append(page.get_text())
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"读取第{i+1}/{total}页...")
        doc.close()
        self._write_text(out, '\n'.join(text))
        return True

    def _pdf_to_image(self, inp, out, cb):
        import fitz
        doc = fitz.open(inp)
        total = len(doc)
        out_dir = os.path.dirname(out)
        base = os.path.splitext(os.path.basename(out))[0]
        ext = os.path.splitext(out)[1]
        for i, page in enumerate(doc):
            if self._cancel: return False
            pix = page.get_pixmap(dpi=200)
            if total == 1:
                pix.save(out)
            else:
                pix.save(os.path.join(out_dir, f"{base}_{i+1}{ext}"))
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"渲染第{i+1}/{total}页...")
        doc.close()
        return True

    def _pdf_to_html(self, inp, out, cb):
        import fitz
        doc = fitz.open(inp)
        if cb: cb(30, "转换中...")
        parts = ["<html><head><meta charset='utf-8'></head><body>"]
        for page in doc:
            parts.append(page.get_text("html"))
        parts.append("</body></html>")
        doc.close()
        self._write_text(out, '\n'.join(parts))
        return True

    def _pdf_to_pptx(self, inp, out, cb):
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        doc = fitz.open(inp)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text = page.get_text().strip()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            if text:
                for line in text.split('\n')[:50]:
                    p = tf.add_paragraph()
                    p.text = line[:200]
                    p.font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = f"(第{i+1}页 - 无文字内容)"
                p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"第{i+1}/{total}页...")
        doc.close()
        prs.save(out)
        return True

    def _pdf_to_xlsx(self, inp, out, cb):
        import fitz
        import openpyxl
        doc = fitz.open(inp)
        wb = openpyxl.Workbook()
        ws = wb.active
        total = len(doc)
        for i, page in enumerate(doc):
            if self._cancel: return False
            text = page.get_text().strip().split('\n')
            for j, line in enumerate(text):
                ws.cell(row=i+1, column=j+1, value=line[:32767])
            if cb:
                cb(20 + int(i * 70 / max(total, 1)), f"第{i+1}/{total}页...")
        doc.close()
        wb.save(out)
        return True

    # ========== Word 转换 ==========

    def _docx_to_pdf(self, inp, out, cb):
        if cb: cb(30, "启动Word...")
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        word = None
        doc = None
        try:
            word = win32com.client.DispatchEx("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            if cb: cb(50, "打开文档...")
            abs_inp = os.path.abspath(inp)
            abs_out = os.path.abspath(out)
            doc = word.Documents.Open(abs_inp, ReadOnly=True)
            if cb: cb(70, "导出PDF...")
            doc.SaveAs2(abs_out, FileFormat=17)
            doc.Close(0)
            doc = None
            word.Quit()
            word = None
            return True
        except Exception as e:
            if doc:
                try: doc.Close(0)
                except: pass
            if word:
                try: word.Quit()
                except: pass
            raise e
        finally:
            pythoncom.CoUninitialize()

    def _docx_to_txt(self, inp, out, cb):
        if cb: cb(30, "读取文档...")
        import docx
        doc = docx.Document(inp)
        text = '\n'.join(p.text for p in doc.paragraphs)
        self._write_text(out, text)
        return True

    def _docx_to_html(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        import docx
        doc = docx.Document(inp)
        body_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                body_parts.append('<br>')
                continue
            tag = 'p'
            if para.style and para.style.name:
                sn = para.style.name.lower()
                if 'heading 1' in sn or '标题 1' in sn: tag = 'h1'
                elif 'heading 2' in sn or '标题 2' in sn: tag = 'h2'
                elif 'heading 3' in sn or '标题 3' in sn: tag = 'h3'
            body_parts.append(f'<{tag}>{self._safe_html(text)}</{tag}>')
        if doc.tables:
            for table in doc.tables:
                body_parts.append('<table border="1" cellpadding="4" style="border-collapse:collapse">')
                for row in table.rows:
                    body_parts.append('<tr>')
                    for cell in row.cells:
                        body_parts.append(f'<td>{self._safe_html(cell.text.strip())}</td>')
                    body_parts.append('</tr>')
                body_parts.append('</table><br>')
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        return True

    def _docx_to_image(self, inp, out, cb):
        if cb: cb(20, "读取文档...")
        import docx
        from PIL import Image, ImageDraw, ImageFont
        doc = docx.Document(inp)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        if not lines:
            lines = ["(空白文档)"]
        font_size = 20
        try:
            font = ImageFont.truetype("msyh.ttc", font_size)
        except Exception:
            font = ImageFont.load_default()
        line_h = font_size + 10
        img_w = 1000
        img_h = max(200, len(lines) * line_h + 100)
        img = Image.new('RGB', (img_w, img_h), 'white')
        draw = ImageDraw.Draw(img)
        y = 20
        for line in lines:
            draw.text((20, y), line[:120], fill='black', font=font)
            y += line_h
        img.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _docx_to_pptx(self, inp, out, cb):
        if cb: cb(20, "读取文档...")
        import docx
        from pptx import Presentation
        from pptx.util import Inches, Pt
        doc = docx.Document(inp)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        chunk_size = 15
        for i in range(0, max(len(lines), 1), chunk_size):
            if self._cancel: return False
            chunk = lines[i:i+chunk_size]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            if chunk:
                for j, line in enumerate(chunk):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line[:200]
                    p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(len(lines), 1)), f"段落{i+1}...")
        prs.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _docx_to_md(self, inp, out, cb):
        if cb: cb(30, "读取文档...")
        import docx
        doc = docx.Document(inp)
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append('')
                continue
            if para.style and para.style.name:
                sn = para.style.name.lower()
                if 'heading 1' in sn or '标题 1' in sn:
                    lines.append(f'# {text}')
                    continue
                elif 'heading 2' in sn or '标题 2' in sn:
                    lines.append(f'## {text}')
                    continue
                elif 'heading 3' in sn or '标题 3' in sn:
                    lines.append(f'### {text}')
                    continue
            lines.append(text)
        self._write_text(out, '\n\n'.join(lines))
        if cb: cb(100, "转换完成")
        return True

    def _doc_copy(self, inp, out, cb):
        if cb: cb(50, "复制转换中...")
        import shutil
        shutil.copy2(inp, out)
        return True

    # ========== Excel (.xlsx) 转换 ==========

    def _xlsx_to_pdf(self, inp, out, cb):
        if cb: cb(30, "读取表格...")
        import openpyxl
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        if cb: cb(50, "生成PDF...")
        data = self._make_table_data(ws)
        pdf_doc = SimpleDocTemplate(out, pagesize=landscape(A4))
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.Color(0.95, 0.95, 0.95)]),
        ]))
        pdf_doc.build([table])
        return True

    def _xlsx_to_csv(self, inp, out, cb):
        import openpyxl
        import csv
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        with open(out, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            total = ws.max_row or 1
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if self._cancel: return False
                writer.writerow(row)
                if cb:
                    cb(20 + int(i * 70 / total), f"写入第{i+1}行...")
        return True

    def _xlsx_to_txt(self, inp, out, cb):
        import openpyxl
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        lines = []
        for row in ws.iter_rows(values_only=True):
            lines.append('\t'.join(str(c) if c is not None else '' for c in row))
        self._write_text(out, '\n'.join(lines))
        return True

    def _xlsx_to_image(self, inp, out, cb):
        if cb: cb(20, "读取表格...")
        import openpyxl
        from PIL import Image, ImageDraw, ImageFont
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        data = self._make_table_data(ws)
        try:
            font = ImageFont.truetype("msyh.ttc", 16)
        except Exception:
            font = ImageFont.load_default()
        cell_w, cell_h = 120, 30
        col_w = [cell_w] * max(len(data[0]) if data else 1, 1)
        for row_idx, row in enumerate(data):
            for col_idx, val in enumerate(row):
                est = len(str(val)) * 10 + 20
                if col_idx < len(col_w) and est > col_w[col_idx]:
                    col_w[col_idx] = min(est, 300)
        img_w = max(400, sum(col_w) + 20)
        img_h = max(200, len(data) * cell_h + 40)
        img = Image.new('RGB', (img_w, img_h), 'white')
        draw = ImageDraw.Draw(img)
        draw.rectangle([0, 0, img_w-1, img_h-1], outline='#cccccc')
        y = 10
        for row_idx, row in enumerate(data):
            x = 10
            if row_idx == 0:
                draw.rectangle([0, y, img_w, y+cell_h], fill='#4472C4')
            for col_idx, val in enumerate(row):
                cw = col_w[col_idx] if col_idx < len(col_w) else cell_w
                fill = '#D6E4F0' if row_idx == 0 else ('white' if row_idx % 2 == 1 else '#F2F2F2')
                draw.rectangle([x, y, x+cw, y+cell_h], fill=fill, outline='#cccccc')
                text_color = 'white' if row_idx == 0 else 'black'
                draw.text((x+4, y+6), str(val)[:30], fill=text_color, font=font)
                x += cw
            y += cell_h
        img.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _xlsx_to_html(self, inp, out, cb):
        if cb: cb(20, "读取表格...")
        import openpyxl
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        rows_html = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            tag = 'th' if i == 0 else 'td'
            cells = ''.join(f'<{tag}>{self._safe_html(str(c) if c is not None else "")}</{tag}>' for c in row)
            rows_html.append(f'<tr>{cells}</tr>')
        body = f'<table border="1" cellpadding="6" style="border-collapse:collapse;font-size:12px">{"".join(rows_html)}</table>'
        self._write_text(out, self._build_html_page(body))
        if cb: cb(100, "转换完成")
        return True

    def _xlsx_to_md(self, inp, out, cb):
        if cb: cb(20, "读取表格...")
        import openpyxl
        wb = openpyxl.load_workbook(inp, data_only=True)
        ws = wb.active
        lines = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            vals = [str(c) if c is not None else '' for c in row]
            lines.append('| ' + ' | '.join(vals) + ' |')
            if i == 0:
                lines.append('| ' + ' | '.join(['---'] * len(vals)) + ' |')
        self._write_text(out, '\n'.join(lines))
        if cb: cb(100, "转换完成")
        return True

    # ========== Excel97 (.xls) 转换 ==========

    def _xls_to_xlsx(self, inp, out, cb):
        if cb: cb(20, "读取Excel97...")
        import xlrd
        import openpyxl
        wb_old = xlrd.open_workbook(inp)
        ws_old = wb_old.sheet_by_index(0)
        wb_new = openpyxl.Workbook()
        ws_new = wb_new.active
        for r in range(ws_old.nrows):
            if self._cancel: return False
            for c in range(ws_old.ncols):
                cell = ws_old.cell(r, c)
                ws_new.cell(row=r+1, column=c+1, value=cell.value)
            if cb and r % 50 == 0:
                cb(20 + int(r * 70 / max(ws_old.nrows, 1)), f"第{r+1}/{ws_old.nrows}行...")
        wb_new.save(out)
        return True

    def _xls_to_pdf(self, inp, out, cb):
        if cb: cb(20, "读取Excel97...")
        import xlrd
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        data = []
        for r in range(ws.nrows):
            row = []
            for c in range(ws.ncols):
                cell = ws.cell(r, c)
                row.append(str(cell.value) if cell.value is not None else '')
            data.append(row)
        if not data:
            data = [['(空表格)']]
        if cb: cb(50, "生成PDF...")
        pdf_doc = SimpleDocTemplate(out, pagesize=landscape(A4))
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.Color(0.95, 0.95, 0.95)]),
        ]))
        pdf_doc.build([table])
        return True

    def _xls_to_csv(self, inp, out, cb):
        import xlrd
        import csv
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        with open(out, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            for r in range(ws.nrows):
                if self._cancel: return False
                row = [str(ws.cell(r, c).value) if ws.cell(r, c).value is not None else '' for c in range(ws.ncols)]
                writer.writerow(row)
                if cb and r % 50 == 0:
                    cb(20 + int(r * 70 / max(ws.nrows, 1)), f"第{r+1}/{ws.nrows}行...")
        return True

    def _xls_to_txt(self, inp, out, cb):
        import xlrd
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        lines = []
        for r in range(ws.nrows):
            row = '\t'.join(str(ws.cell(r, c).value) if ws.cell(r, c).value is not None else '' for c in range(ws.ncols))
            lines.append(row)
        self._write_text(out, '\n'.join(lines))
        return True

    def _xls_to_image(self, inp, out, cb):
        if cb: cb(20, "读取Excel97...")
        import xlrd
        from PIL import Image, ImageDraw, ImageFont
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        data = []
        for r in range(ws.nrows):
            row = [str(ws.cell(r, c).value) if ws.cell(r, c).value is not None else '' for c in range(ws.ncols)]
            data.append(row)
        if not data:
            data = [['(空表格)']]
        try:
            font = ImageFont.truetype("msyh.ttc", 16)
        except Exception:
            font = ImageFont.load_default()
        cell_w, cell_h = 120, 30
        col_w = [cell_w] * max(len(data[0]), 1)
        for row_idx, row in enumerate(data):
            for col_idx, val in enumerate(row):
                est = len(str(val)) * 10 + 20
                if col_idx < len(col_w) and est > col_w[col_idx]:
                    col_w[col_idx] = min(est, 300)
        img_w = max(400, sum(col_w) + 20)
        img_h = max(200, len(data) * cell_h + 40)
        img = Image.new('RGB', (img_w, img_h), 'white')
        draw = ImageDraw.Draw(img)
        y = 10
        for row_idx, row in enumerate(data):
            x = 10
            if row_idx == 0:
                draw.rectangle([0, y, img_w, y+cell_h], fill='#4472C4')
            for col_idx, val in enumerate(row):
                cw = col_w[col_idx] if col_idx < len(col_w) else cell_w
                fill = '#D6E4F0' if row_idx == 0 else ('white' if row_idx % 2 == 1 else '#F2F2F2')
                draw.rectangle([x, y, x+cw, y+cell_h], fill=fill, outline='#cccccc')
                text_color = 'white' if row_idx == 0 else 'black'
                draw.text((x+4, y+6), str(val)[:30], fill=text_color, font=font)
                x += cw
            y += cell_h
        img.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _xls_to_html(self, inp, out, cb):
        if cb: cb(20, "读取Excel97...")
        import xlrd
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        rows_html = []
        for r in range(ws.nrows):
            cells = ''.join(f'<{"th" if r==0 else "td"}>{self._safe_html(str(ws.cell(r,c).value) if ws.cell(r,c).value is not None else "")}</{"th" if r==0 else "td"}>' for c in range(ws.ncols))
            rows_html.append(f'<tr>{cells}</tr>')
        body = f'<table border="1" cellpadding="6" style="border-collapse:collapse;font-size:12px">{"".join(rows_html)}</table>'
        self._write_text(out, self._build_html_page(body))
        if cb: cb(100, "转换完成")
        return True

    def _xls_to_md(self, inp, out, cb):
        if cb: cb(20, "读取Excel97...")
        import xlrd
        wb = xlrd.open_workbook(inp)
        ws = wb.sheet_by_index(0)
        lines = []
        for r in range(ws.nrows):
            vals = [str(ws.cell(r, c).value) if ws.cell(r, c).value is not None else '' for c in range(ws.ncols)]
            lines.append('| ' + ' | '.join(vals) + ' |')
            if r == 0:
                lines.append('| ' + ' | '.join(['---'] * len(vals)) + ' |')
        self._write_text(out, '\n'.join(lines))
        return True

    # ========== CSV 转换 ==========

    def _csv_to_pdf(self, inp, out, cb):
        if cb: cb(30, "读取CSV...")
        import csv
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        data = []
        with open(inp, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for row in reader:
                data.append(row)
        if not data:
            data = [['(空文件)']]
        if cb: cb(50, "生成PDF...")
        pdf_doc = SimpleDocTemplate(out, pagesize=landscape(A4))
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ]))
        pdf_doc.build([table])
        return True

    def _csv_to_html(self, inp, out, cb):
        if cb: cb(20, "读取CSV...")
        import csv
        rows_html = []
        with open(inp, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                cells = ''.join(f'<{"th" if i==0 else "td"}>{self._safe_html(c)}</{"th" if i==0 else "td"}>' for c in row)
                rows_html.append(f'<tr>{cells}</tr>')
        body = f'<table border="1" cellpadding="6" style="border-collapse:collapse">{"".join(rows_html)}</table>'
        self._write_text(out, self._build_html_page(body))
        return True

    def _csv_to_xlsx(self, inp, out, cb):
        if cb: cb(20, "读取CSV...")
        import csv, openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(inp, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if self._cancel: return False
                ws.append(row)
                if cb and i % 100 == 0:
                    cb(20 + min(70, i // 10), f"写入第{i+1}行...")
        wb.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _csv_to_md(self, inp, out, cb):
        if cb: cb(20, "读取CSV...")
        import csv
        lines = []
        with open(inp, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                vals = [c for c in row]
                lines.append('| ' + ' | '.join(vals) + ' |')
                if i == 0:
                    lines.append('| ' + ' | '.join(['---'] * len(vals)) + ' |')
        self._write_text(out, '\n'.join(lines))
        return True

    # ========== TXT 转换 ==========

    def _txt_to_xlsx(self, inp, out, cb):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(inp, 'r', encoding='utf-8') as f:
            for i, line in enumerate(f):
                if self._cancel: return False
                parts = line.strip().split('\t')
                if len(parts) == 1:
                    parts = line.strip().split(',')
                ws.append(parts)
                if cb and i % 100 == 0:
                    cb(20 + min(70, i // 10), f"写入第{i+1}行...")
        wb.save(out)
        return True

    def _txt_to_docx(self, inp, out, cb):
        if cb: cb(30, "生成Word...")
        import docx
        doc = docx.Document()
        with open(inp, 'r', encoding='utf-8') as f:
            for line in f:
                if self._cancel: return False
                text = line.rstrip()
                if text.startswith('# '):
                    doc.add_heading(text[2:], level=1)
                elif text.startswith('## '):
                    doc.add_heading(text[3:], level=2)
                elif text.startswith('### '):
                    doc.add_heading(text[4:], level=3)
                elif text.strip() == '':
                    doc.add_paragraph('')
                else:
                    doc.add_paragraph(text)
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_pptx(self, inp, out, cb):
        if cb: cb(20, "生成PPT...")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        with open(inp, 'r', encoding='utf-8') as f:
            lines = [l.rstrip() for l in f if l.strip()]
        chunk_size = 12
        for i in range(0, max(len(lines), 1), chunk_size):
            if self._cancel: return False
            chunk = lines[i:i+chunk_size]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, line in enumerate(chunk):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if line.startswith('# '):
                    p.text = line[2:]
                    p.font.size = Pt(28)
                    p.font.bold = True
                else:
                    p.text = line[:200]
                    p.font.size = Pt(18)
            if cb:
                cb(20 + int(i * 70 / max(len(lines), 1)), f"生成幻灯片...")
        prs.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_pdf(self, inp, out, cb):
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        font_name = 'SimSun'
        try:
            for fp in ["C:/Windows/Fonts/simsun.ttc", "C:/Windows/Fonts/msyh.ttc"]:
                if os.path.exists(fp):
                    pdfmetrics.registerFont(TTFont(font_name, fp))
                    break
            else:
                font_name = 'Helvetica'
        except Exception:
            font_name = 'Helvetica'
        if cb: cb(30, "读取文本...")
        with open(inp, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        pdf_doc = SimpleDocTemplate(out, pagesize=A4)
        styles = getSampleStyleSheet()
        style = styles['Normal']
        style.fontName = font_name
        story = []
        for line in lines:
            safe = line.rstrip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            if safe:
                story.append(Paragraph(safe, style))
            else:
                story.append(Spacer(1, 12))
        if cb: cb(70, "生成PDF...")
        pdf_doc.build(story)
        return True

    def _txt_to_html(self, inp, out, cb):
        if cb: cb(20, "转换中...")
        with open(inp, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        body_parts = []
        for line in lines:
            text = line.rstrip()
            if not text:
                body_parts.append('<br>')
            elif text.startswith('# '):
                body_parts.append(f'<h1>{self._safe_html(text[2:])}</h1>')
            elif text.startswith('## '):
                body_parts.append(f'<h2>{self._safe_html(text[3:])}</h2>')
            elif text.startswith('### '):
                body_parts.append(f'<h3>{self._safe_html(text[4:])}</h3>')
            else:
                body_parts.append(f'<p>{self._safe_html(text)}</p>')
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        if cb: cb(100, "转换完成")
        return True

    def _txt_to_md(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        import shutil
        shutil.copy2(inp, out)
        if cb: cb(100, "转换完成")
        return True

    # ========== PPT 转换 ==========

    def _pptx_to_pdf(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        prs = Presentation(inp)
        if cb: cb(50, "生成PDF...")
        pdf_doc = SimpleDocTemplate(out, pagesize=landscape(A4))
        styles = getSampleStyleSheet()
        story = []
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            story.append(Paragraph(f"--- 幻灯片 {i+1} ---", styles['Heading1']))
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    safe = self._safe_html(shape.text)
                    story.append(Paragraph(safe, styles['Normal']))
            story.append(Spacer(1, 24))
            if cb:
                cb(50 + int(i * 40 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        pdf_doc.build(story)
        return True

    def _pptx_to_txt(self, inp, out, cb):
        from pptx import Presentation
        prs = Presentation(inp)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
        self._write_text(out, '\n\n'.join(text_parts))
        if cb: cb(100, "转换完成")
        return True

    def _pptx_to_image(self, inp, out, cb):
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        prs = Presentation(inp)
        out_dir = os.path.dirname(out)
        base = os.path.splitext(os.path.basename(out))[0]
        ext = os.path.splitext(out)[1]
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            img = Image.new('RGB', (1280, 720), 'white')
            draw = ImageDraw.Draw(img)
            y = 50
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    for line in shape.text.split('\n'):
                        if line.strip():
                            try:
                                font = ImageFont.truetype("msyh.ttc", 18)
                            except Exception:
                                font = ImageFont.load_default()
                            draw.text((60, y), line.strip(), fill='black', font=font)
                            y += 30
                            if y > 680: break
            if total == 1:
                img.save(out)
            else:
                img.save(os.path.join(out_dir, f"{base}_{i+1}{ext}"))
            if cb:
                cb(20 + int(i * 70 / total), f"幻灯片{i+1}/{total}...")
        return True

    def _pptx_to_docx(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        import docx
        prs = Presentation(inp)
        doc = docx.Document()
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            doc.add_heading(f'幻灯片 {i+1}', level=2)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    doc.add_paragraph(shape.text.strip())
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _pptx_to_html(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        prs = Presentation(inp)
        body_parts = []
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            body_parts.append(f'<h2>幻灯片 {i+1}</h2>')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    body_parts.append(f'<p>{self._safe_html(shape.text.strip())}</p>')
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        self._write_text(out, self._build_html_page('\n'.join(body_parts)))
        return True

    def _pptx_to_md(self, inp, out, cb):
        if cb: cb(30, "读取演示文稿...")
        from pptx import Presentation
        prs = Presentation(inp)
        lines = []
        for i, slide in enumerate(prs.slides):
            if self._cancel: return False
            lines.append(f'## 幻灯片 {i+1}')
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip())
            lines.append('')
            if cb:
                cb(20 + int(i * 70 / max(len(prs.slides), 1)), f"幻灯片{i+1}...")
        self._write_text(out, '\n'.join(lines))
        return True

    def _ppt_to_pdf(self, inp, out, cb):
        if cb: cb(30, "启动PowerPoint...")
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        ppt = None
        pres = None
        try:
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            ppt.Visible = False
            if cb: cb(50, "打开演示文稿...")
            abs_inp = os.path.abspath(inp)
            abs_out = os.path.abspath(out)
            pres = ppt.Presentations.Open(abs_inp, WithWindow=False)
            if cb: cb(70, "导出PDF...")
            pres.SaveAs(abs_out, 32)
            pres.Close()
            pres = None
            ppt.Quit()
            ppt = None
            return True
        except Exception as e:
            if pres:
                try: pres.Close()
                except: pass
            if ppt:
                try: ppt.Quit()
                except: pass
            raise e
        finally:
            pythoncom.CoUninitialize()

    def _ppt_to_txt(self, inp, out, cb):
        if cb: cb(30, "启动PowerPoint...")
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        ppt = None
        pres = None
        try:
            ppt = win32com.client.DispatchEx("PowerPoint.Application")
            ppt.Visible = False
            if cb: cb(50, "读取演示文稿...")
            abs_inp = os.path.abspath(inp)
            pres = ppt.Presentations.Open(abs_inp, WithWindow=False)
            texts = []
            for slide in pres.Slides:
                if self._cancel: return False
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        texts.append(shape.TextFrame.TextRange.Text)
            pres.Close()
            pres = None
            ppt.Quit()
            ppt = None
            self._write_text(out, '\n\n'.join(texts))
            if cb: cb(100, "转换完成")
            return True
        except Exception as e:
            if pres:
                try: pres.Close()
                except: pass
            if ppt:
                try: ppt.Quit()
                except: pass
            raise e
        finally:
            pythoncom.CoUninitialize()

    # ========== 图片 转换 ==========

    def _image_to_pdf(self, inp, out, cb):
        from PIL import Image
        if cb: cb(30, "处理图片...")
        img = Image.open(inp)
        if img.mode == 'RGBA':
            bg = Image.new('RGB', img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            img = bg
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        if cb: cb(70, "保存PDF...")
        img.save(out, 'PDF')
        return True

    def _image_to_docx(self, inp, out, cb):
        if cb: cb(30, "生成Word...")
        import docx
        doc = docx.Document()
        doc.add_picture(inp)
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== HTML 转换 ==========

    def _html_to_pdf(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', ' ', content)
        text = re.sub(r'\s+', ' ', text).strip()
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try:
            os.remove(temp)
        except OSError:
            pass
        return result

    def _html_to_docx(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        import docx
        content = self._read_text(inp)
        doc = docx.Document()
        for match in re.finditer(r'<h([1-3])[^>]*>(.*?)</h\1>', content, re.DOTALL):
            doc.add_heading(re.sub(r'<[^>]+>', '', match.group(2)).strip(), level=int(match.group(1)))
        for match in re.finditer(r'<p[^>]*>(.*?)</p>', content, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text:
                doc.add_paragraph(text)
        if not doc.paragraphs:
            text = re.sub(r'<[^>]+>', '', content)
            for line in text.strip().split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_txt(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', '', content)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_md(self, inp, out, cb):
        if cb: cb(30, "转换中...")
        content = self._read_text(inp)
        text = re.sub(r'<[^>]+>', '', content)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _html_to_xlsx(self, inp, out, cb):
        if cb: cb(30, "解析HTML...")
        import openpyxl
        content = self._read_text(inp)
        wb = openpyxl.Workbook()
        ws = wb.active
        tables = re.findall(r'<table[^>]*>(.*?)</table>', content, re.DOTALL)
        if tables:
            for table_html in tables:
                rows = re.findall(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL)
                for ri, row_html in enumerate(rows):
                    cells = re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', row_html, re.DOTALL)
                    for ci, cell in enumerate(cells):
                        ws.cell(row=ri+1, column=ci+1, value=re.sub(r'<[^>]+>', '', cell).strip())
                break
        else:
            text = re.sub(r'<[^>]+>', '', content)
            for i, line in enumerate(text.strip().split('\n')[:100]):
                if line.strip():
                    ws.cell(row=i+1, column=1, value=line.strip())
        wb.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== Markdown 转换 ==========

    def _md_to_html(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        self._write_text(out, self._build_html_page(html))
        if cb: cb(100, "转换完成")
        return True

    def _md_to_pdf(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        temp_html = out + '.tmp.html'
        self._write_text(temp_html, self._build_html_page(html))
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'\s+', ' ', text).strip()
        temp_txt = out + '.tmp.txt'
        self._write_text(temp_txt, text)
        result = self._txt_to_pdf(temp_txt, out, cb)
        try:
            os.remove(temp_html)
            os.remove(temp_txt)
        except OSError:
            pass
        return result

    def _md_to_docx(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        import docx
        md = MarkdownIt()
        tokens = md.parse(self._read_text(inp))
        doc = docx.Document()
        for token in tokens:
            if token.type == 'heading_open':
                level = int(token.tag[1])
            elif token.type == 'inline' and token.content.strip():
                if any(t.type == 'heading_open' for t in tokens[:tokens.index(token)+1] if t.type in ('heading_open','inline','paragraph_open')):
                    pass
        html = md.render(self._read_text(inp))
        for match in re.finditer(r'<h([1-3])[^>]*>(.*?)</h\1>', html, re.DOTALL):
            doc.add_heading(re.sub(r'<[^>]+>', '', match.group(2)).strip(), level=int(match.group(1)))
        for match in re.finditer(r'<li>(.*?)</li>', html, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text:
                doc.add_paragraph(text, style='List Bullet')
        for match in re.finditer(r'<p>(.*?)</p>', html, re.DOTALL):
            text = re.sub(r'<[^>]+>', '', match.group(1)).strip()
            if text and not any(doc.paragraphs[-1].text == text for _ in [0] if doc.paragraphs):
                doc.add_paragraph(text)
        if not doc.paragraphs:
            text = re.sub(r'<[^>]+>', '', html)
            for line in text.strip().split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _md_to_txt(self, inp, out, cb):
        if cb: cb(20, "解析Markdown...")
        from markdown_it import MarkdownIt
        md = MarkdownIt()
        html = md.render(self._read_text(inp))
        text = re.sub(r'<[^>]+>', '', html)
        text = re.sub(r'\s+', ' ', text).strip()
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    # ========== EPUB 转换 ==========

    def _epub_to_text(self, inp):
        from ebooklib import epub
        book = epub.read_epub(inp)
        texts = []
        for item in book.get_items():
            if item.get_type() == 9:
                content = item.get_content().decode('utf-8', errors='replace')
                texts.append(re.sub(r'<[^>]+>', ' ', content))
        return '\n\n'.join(texts)

    def _epub_to_pdf(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        text = self._epub_to_text(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except: pass
        return result

    def _epub_to_txt(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        text = self._epub_to_text(inp)
        self._write_text(out, text)
        if cb: cb(100, "转换完成")
        return True

    def _epub_to_html(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        from ebooklib import epub
        book = epub.read_epub(inp)
        body = []
        for item in book.get_items():
            if item.get_type() == 9:
                body.append(item.get_content().decode('utf-8', errors='replace'))
        self._write_text(out, self._build_html_page('\n'.join(body)))
        if cb: cb(100, "转换完成")
        return True

    def _epub_to_docx(self, inp, out, cb):
        if cb: cb(20, "读取EPUB...")
        import docx
        text = self._epub_to_text(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== RTF 转换 ==========

    def _rtf_read(self, inp):
        from striprtf.striprtf import rtf_to_text
        text = self._read_text(inp)
        return rtf_to_text(text)

    def _rtf_to_txt(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        self._write_text(out, self._rtf_read(inp))
        if cb: cb(100, "转换完成")
        return True

    def _rtf_to_pdf(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        text = self._rtf_read(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except: pass
        return result

    def _rtf_to_docx(self, inp, out, cb):
        if cb: cb(30, "解析RTF...")
        import docx
        text = self._rtf_read(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    # ========== ODT 转换 ==========

    def _odt_read_text(self, inp):
        from odf.opendocument import load
        from odf.text import P
        doc = load(inp)
        texts = []
        for p in doc.getElementsByType(P):
            texts.append(str(p))
        return '\n'.join(texts)

    def _odt_to_pdf(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        text = self._odt_read_text(inp)
        temp = out + '.tmp.txt'
        self._write_text(temp, text)
        result = self._txt_to_pdf(temp, out, cb)
        try: os.remove(temp)
        except: pass
        return result

    def _odt_to_docx(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        import docx
        text = self._odt_read_text(inp)
        doc = docx.Document()
        for line in text.split('\n'):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(out)
        if cb: cb(100, "转换完成")
        return True

    def _odt_to_txt(self, inp, out, cb):
        if cb: cb(30, "读取ODT...")
        self._write_text(out, self._odt_read_text(inp))
        if cb: cb(100, "转换完成")
        return True